"""
Export PowerPoint (.pptx) và PDF có brand công ty.

- PPTX: dùng python-pptx (header logo/màu brand, slide per report, chart PNG)
- PDF:  dùng weasyprint nếu có, fallback LibreOffice headless, fallback base64 HTML

Cài thêm nếu cần:
  pip install python-pptx weasyprint
"""

from __future__ import annotations

import io
import logging
import re
from typing import Any

_logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Branding helpers
# ---------------------------------------------------------------------------

def _default_branding() -> dict[str, str]:
    import os
    return {
        "product_name": os.getenv("BRAND_PRODUCT_NAME", "AI BI Smart"),
        "primary_color": os.getenv("BRAND_PRIMARY_COLOR", "#0f766e"),
        "logo_url": os.getenv("BRAND_LOGO_URL", ""),
    }


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    h = (hex_color or "#0f766e").lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    r = int(h[0:2], 16)
    g = int(h[2:4], 16)
    b = int(h[4:6], 16)
    return r, g, b


def _clean_md(text: str) -> str:
    """Bỏ markdown cơ bản để dán vào slide/PDF."""
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text or "")
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return text.strip()


# ---------------------------------------------------------------------------
# PPTX export
# ---------------------------------------------------------------------------

def create_pptx_report(
    reports: list[dict[str, Any]],
    *,
    title: str = "Báo cáo BI",
    branding: dict[str, str] | None = None,
    chart_images: list[bytes | None] | None = None,
) -> bytes:
    """
    Tạo file PPTX từ danh sách report.

    reports: list[{query, insight, data, chart_type}]
    chart_images: list[PNG bytes hoặc None] tương ứng với reports.

    Returns bytes (.pptx)
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Emu, Inches, Pt
    except ImportError as exc:
        raise ValueError(
            "Thiếu python-pptx. Cài: pip install python-pptx"
        ) from exc

    brand = branding or _default_branding()
    product_name = brand.get("product_name", "AI BI Smart")
    primary_hex = brand.get("primary_color", "#0f766e")
    r, g, b = _hex_to_rgb(primary_hex)
    brand_color = RGBColor(r, g, b)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank

    # --- Slide tiêu đề ---
    slide_title = prs.slides.add_slide(blank_layout)
    _add_brand_header(slide_title, product_name, brand_color, prs)

    # Tiêu đề chính
    txBox = slide_title.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = brand_color

    # Subtitle
    txBox2 = slide_title.shapes.add_textbox(
        Inches(0.5), Inches(4.2), Inches(12.33), Inches(0.6)
    )
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = product_name
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0] if p2.runs else p2.add_run()
    run2.font.size = Pt(16)
    run2.font.color.rgb = RGBColor(100, 100, 100)

    # --- Slide mỗi report ---
    imgs = chart_images or []
    for i, rep in enumerate(reports):
        slide = prs.slides.add_slide(blank_layout)
        _add_brand_header(slide, product_name, brand_color, prs)

        q = str(rep.get("query") or f"Báo cáo {i + 1}")
        insight = _clean_md(str(rep.get("insight") or ""))

        # Tiêu đề query
        txQ = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.7))
        tfQ = txQ.text_frame
        pQ = tfQ.paragraphs[0]
        pQ.text = q[:120]
        runQ = pQ.runs[0] if pQ.runs else pQ.add_run()
        runQ.font.size = Pt(20)
        runQ.font.bold = True
        runQ.font.color.rgb = brand_color

        # Chart image nếu có
        img_bytes: bytes | None = imgs[i] if i < len(imgs) else None
        if img_bytes:
            try:
                slide.shapes.add_picture(
                    io.BytesIO(img_bytes),
                    Inches(0.4), Inches(1.7), Inches(7.5), Inches(4.5),
                )
                # Insight bên phải
                txI = slide.shapes.add_textbox(Inches(8.2), Inches(1.7), Inches(4.8), Inches(4.5))
                tfI = txI.text_frame
                tfI.word_wrap = True
                _fill_insight_tf(tfI, insight)
            except Exception:
                # Fallback: insight full width
                _add_full_insight(slide, insight)
        else:
            _add_full_insight(slide, insight)

        # Bảng dữ liệu nhỏ (5 dòng đầu)
        data_rows = (rep.get("data") or [])[:5]
        if data_rows and not img_bytes:
            _add_data_table(slide, data_rows, brand_color, prs)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_brand_header(slide: Any, product_name: str, brand_color: Any, prs: Any) -> None:
    """Thêm dải header màu brand + tên sản phẩm."""
    try:
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.oxml.ns import qn
    except ImportError:
        return
    # Header rectangle
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), prs.slide_width, Inches(0.75),
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = brand_color
    shape.line.fill.background()

    # Product name
    tf = shape.text_frame
    tf.text = product_name
    p = tf.paragraphs[0]
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = __import__("pptx").dml.color.RGBColor(255, 255, 255)


def _fill_insight_tf(tf: Any, insight: str) -> None:
    try:
        from pptx.util import Pt
    except ImportError:
        return
    tf.word_wrap = True
    paras = insight.split("\n")[:15]
    for j, line in enumerate(paras):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line[:200]
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(11)


def _add_full_insight(slide: Any, insight: str) -> None:
    try:
        from pptx.util import Inches, Pt
    except ImportError:
        return
    txI = slide.shapes.add_textbox(Inches(0.4), Inches(1.7), Inches(12.5), Inches(5.2))
    tfI = txI.text_frame
    _fill_insight_tf(tfI, insight)


def _add_data_table(slide: Any, rows: list[dict[str, Any]], brand_color: Any, prs: Any) -> None:
    if not rows:
        return
    try:
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return

    cols = list(rows[0].keys())[:6]
    n_rows = len(rows) + 1  # +1 header
    n_cols = len(cols)

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.5),
    ).table

    # Header
    for j, col in enumerate(cols):
        cell = tbl.cell(0, j)
        cell.text = str(col)[:24]
        cell.fill.solid()
        cell.fill.fore_color.rgb = brand_color
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)

    # Dữ liệu
    for i, row in enumerate(rows):
        for j, col in enumerate(cols):
            cell = tbl.cell(i + 1, j)
            v = row.get(col, "")
            cell.text = str(v)[:30] if v is not None else ""
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(9)


# ---------------------------------------------------------------------------
# PDF export
# ---------------------------------------------------------------------------

def create_pdf_report(
    reports: list[dict[str, Any]],
    *,
    title: str = "Báo cáo BI",
    branding: dict[str, str] | None = None,
    chart_images: list[bytes | None] | None = None,
) -> bytes:
    """
    Tạo PDF từ danh sách report.
    Chiến lược: weasyprint (HTML→PDF) → PPTX→base64 fallback.
    """
    brand = branding or _default_branding()
    product_name = brand.get("product_name", "AI BI Smart")
    primary_color = brand.get("primary_color", "#0f766e")

    html = _build_html(reports, title=title, product_name=product_name,
                        primary_color=primary_color, chart_images=chart_images or [])

    # Thử weasyprint
    try:
        from weasyprint import HTML  # type: ignore[import]
        return HTML(string=html).write_pdf()
    except ImportError:
        pass
    except Exception as exc:
        _logger.warning("weasyprint failed: %s", exc)

    # Fallback: trả HTML bytes (browser in thành PDF)
    _logger.info("PDF fallback: trả HTML — dùng browser Print → Save as PDF")
    return html.encode("utf-8")


def _build_html(
    reports: list[dict[str, Any]],
    *,
    title: str,
    product_name: str,
    primary_color: str,
    chart_images: list[bytes | None],
) -> str:
    import base64

    def _img_tag(img_bytes: bytes | None) -> str:
        if not img_bytes:
            return ""
        b64 = base64.b64encode(img_bytes).decode("ascii")
        return f'<img src="data:image/png;base64,{b64}" style="max-width:100%;margin:12px 0;" />'

    sections = ""
    for i, rep in enumerate(reports):
        q = rep.get("query", f"Báo cáo {i + 1}")
        insight = _clean_md(str(rep.get("insight") or ""))
        img_bytes = chart_images[i] if i < len(chart_images) else None

        # Bảng dữ liệu
        data_rows = (rep.get("data") or [])[:20]
        table_html = ""
        if data_rows:
            cols = list(data_rows[0].keys())[:8]
            ths = "".join(f"<th>{c}</th>" for c in cols)
            trs = ""
            for row in data_rows:
                tds = "".join(f"<td>{row.get(c, '')}</td>" for c in cols)
                trs += f"<tr>{tds}</tr>"
            table_html = f"<table><thead><tr>{ths}</tr></thead><tbody>{trs}</tbody></table>"

        sections += f"""
<div class="report-section">
  <h2>{q}</h2>
  {_img_tag(img_bytes)}
  <p class="insight">{insight.replace(chr(10), "<br/>")}</p>
  {table_html}
</div>
"""

    css = f"""
body {{ font-family: 'Segoe UI', Arial, sans-serif; margin: 0; color: #1a1a1a; }}
header {{ background: {primary_color}; color: white; padding: 18px 32px; }}
header h1 {{ margin: 0; font-size: 22px; }}
header small {{ opacity: 0.8; font-size: 13px; }}
.container {{ max-width: 900px; margin: 0 auto; padding: 24px; }}
.report-section {{ border: 1px solid #e0e0e0; border-radius: 8px; padding: 20px; margin-bottom: 28px; }}
h2 {{ color: {primary_color}; font-size: 17px; margin-top: 0; }}
.insight {{ color: #444; line-height: 1.6; font-size: 14px; }}
table {{ width: 100%; border-collapse: collapse; margin-top: 16px; font-size: 12px; }}
th {{ background: {primary_color}; color: white; padding: 6px 10px; text-align: left; }}
td {{ padding: 5px 10px; border-bottom: 1px solid #eee; }}
@media print {{ .report-section {{ page-break-inside: avoid; }} }}
"""

    return f"""<!DOCTYPE html>
<html lang="vi">
<head><meta charset="utf-8"><title>{title}</title>
<style>{css}</style></head>
<body>
<header>
  <h1>{product_name}</h1>
  <small>{title}</small>
</header>
<div class="container">
{sections}
</div>
</body></html>"""
